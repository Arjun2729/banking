from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Deck settings
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Theme
COLORS = {
    "bg": RGBColor(246, 248, 252),
    "primary": RGBColor(18, 74, 158),
    "secondary": RGBColor(26, 32, 44),
    "muted": RGBColor(84, 96, 112),
    "accent": RGBColor(0, 168, 132),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(220, 226, 236),
}


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg"]
    bg.line.fill.background()


def add_top_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.8), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["secondary"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.82), Inches(1.45), Inches(11.5), Inches(0.65))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Calibri"
        sp.font.size = Pt(18)
        sp.font.color.rgb = COLORS["muted"]


def add_bullets(slide, x, y, w, h, bullets, font_size=22, level2_size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for item in bullets:
        level = 0
        text = item
        if isinstance(item, tuple):
            level, text = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = text
        p.level = level
        p.font.name = "Calibri"
        p.font.size = Pt(font_size if level == 0 else level2_size)
        p.font.color.rgb = COLORS["secondary"] if level == 0 else COLORS["muted"]
        if level == 0:
            p.space_after = Pt(10)


def add_card(slide, x, y, w, h, title, body_lines):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["card"]
    card.line.color.rgb = COLORS["line"]

    title_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(w - 0.5), Inches(0.45))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["primary"]

    body_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.75), Inches(w - 0.5), Inches(h - 0.95))
    btf = body_box.text_frame
    btf.word_wrap = True
    btf.clear()

    for i, line in enumerate(body_lines):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["secondary"]
        p.space_after = Pt(7)


def add_footer(slide, text="NexaBank Project Draft | March 2026"):
    foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(12.0), Inches(0.3))
    tf = foot.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["muted"]
    p.alignment = PP_ALIGN.RIGHT


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "NexaBank: Banking Management System", "Project presentation draft based on the current Flask + MySQL implementation")
add_bullets(slide, 0.9, 2.5, 12.0, 2.7, [
    "A modern web application for account management and day-to-day banking actions",
    "Supports registration, authentication, deposits, withdrawals, transfers, and transaction history",
    "Built with pragmatic architecture for clarity, correctness, and extensibility",
], font_size=24)
add_footer(slide)

# Slide 2: Problem & Vision
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Problem Statement and Vision")
add_card(slide, 0.8, 2.0, 6.1, 4.4, "Problem", [
    "Personal banking actions are often fragmented across multiple tools.",
    "Users need quick, reliable access to balances and transaction records.",
    "Financial operations must be accurate, traceable, and secure.",
])
add_card(slide, 6.9, 2.0, 5.6, 4.4, "Vision", [
    "Create one intuitive dashboard for core banking operations.",
    "Guarantee transactional consistency for money movement.",
    "Deliver a clean UX that feels modern and trustworthy.",
])
add_footer(slide)

# Slide 3: Project Scope
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Project Scope and Core Modules")
add_bullets(slide, 0.9, 2.0, 12.0, 4.8, [
    "Authentication and User Management",
    (1, "Register, login, logout, password hashing, profile updates"),
    "Account Lifecycle",
    (1, "Auto-create checking account on registration"),
    (1, "Create additional savings/checking account (one per type)"),
    "Money Movement",
    (1, "Deposit and withdrawal with balance checks"),
    (1, "Internal transfer with lock ordering to prevent deadlocks"),
    "Insights and Reporting",
    (1, "Dashboard KPIs, recent activity, and full transaction history with filters"),
], font_size=24, level2_size=18)
add_footer(slide)

# Slide 4: Product Walkthrough
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "User Journey Walkthrough")
add_card(slide, 0.8, 2.0, 3.0, 4.5, "1. Onboarding", [
    "Register account",
    "Secure login",
    "Session starts",
])
add_card(slide, 3.95, 2.0, 3.0, 4.5, "2. Daily Banking", [
    "View dashboard",
    "Check balances",
    "Open savings/checking",
])
add_card(slide, 7.1, 2.0, 3.0, 4.5, "3. Transactions", [
    "Deposit / Withdraw",
    "Transfer between accounts",
    "See instant updates",
])
add_card(slide, 10.25, 2.0, 2.3, 4.5, "4. Tracking", [
    "Filter history",
    "Paginated logs",
    "Review details",
])
add_footer(slide)

# Slide 5: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "System Architecture")
add_bullets(slide, 0.9, 2.0, 12.0, 1.1, [
    "Monolithic Flask application with clear separation across routing, DB access, templates, and static assets"
], font_size=19)
add_card(slide, 0.9, 3.0, 3.9, 3.3, "Frontend", [
    "Jinja templates",
    "Reusable base layout",
    "CSS + lightweight JS",
])
add_card(slide, 4.95, 3.0, 3.9, 3.3, "Backend", [
    "Flask routes",
    "Session auth",
    "Validation + business rules",
])
add_card(slide, 9.0, 3.0, 3.4, 3.3, "Database", [
    "MySQL",
    "users / accounts / transactions",
    "Atomic updates with transactions",
])
add_footer(slide)

# Slide 6: Database Design
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Database Schema Design")
add_card(slide, 0.8, 2.0, 4.1, 4.6, "users", [
    "id (PK)",
    "full_name, email (unique), phone",
    "password_hash",
    "created_at",
])
add_card(slide, 5.1, 2.0, 4.1, 4.6, "accounts", [
    "id (PK), user_id (FK)",
    "account_type: checking/savings",
    "account_no (unique)",
    "balance, created_at",
])
add_card(slide, 9.4, 2.0, 3.1, 4.6, "transactions", [
    "account_id (FK)",
    "type, amount",
    "balance_after",
    "description, related_account",
    "created_at",
])
add_footer(slide)

# Slide 7: Transaction Integrity
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Transaction Integrity and Concurrency")
add_bullets(slide, 0.9, 2.0, 12.0, 4.8, [
    "All money mutations are wrapped in explicit MySQL transactions (commit/rollback).",
    "`SELECT ... FOR UPDATE` is used for balance reads before updates.",
    "Transfer flow locks accounts in sorted ID order to avoid deadlocks.",
    "Every operation writes an immutable transaction record with resulting balance.",
    "Validation blocks invalid inputs: negative amounts, same-account transfers, insufficient funds.",
], font_size=21)
add_footer(slide)

# Slide 8: Security and Controls
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Security and Reliability Controls")
add_card(slide, 0.8, 2.0, 6.0, 4.6, "Implemented", [
    "Passwords stored with Werkzeug hashing",
    "Session-based authentication + route guards",
    "Parameterized SQL queries to reduce injection risk",
    "Scoped account access by authenticated user ID",
    "Centralized error handling with user-safe messaging",
])
add_card(slide, 6.95, 2.0, 5.6, 4.6, "Next Hardening Steps", [
    "Enable CSRF protection for all forms",
    "Move secrets to secure environment management",
    "Add rate limiting for login and sensitive endpoints",
    "Introduce audit logs and alerting for anomalous activity",
])
add_footer(slide)

# Slide 9: UI and UX
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "UI/UX Highlights")
add_bullets(slide, 0.9, 2.0, 12.0, 4.8, [
    "Consistent component system across templates (`base.html`, form cards, tables, badges).",
    "Dashboard-first experience with total balance, account list, and recent transaction view.",
    "Dedicated pages for deposit, withdrawal, transfer, account detail, and profile settings.",
    "Transaction history supports account/type/date filtering with pagination.",
    "Visual cues for credits vs debits improve scanning and reduce user error.",
], font_size=21)
add_footer(slide)

# Slide 10: Setup and Demo Plan
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Local Setup and Demo Flow")
add_card(slide, 0.8, 2.1, 6.0, 4.5, "How to Run", [
    "1) Install dependencies from requirements.txt",
    "2) Provision MySQL and run schema.sql",
    "3) Set DB_HOST / DB_USER / DB_PASSWORD / DB_NAME",
    "4) Start with: python app.py",
    "5) Open http://localhost:5000",
])
add_card(slide, 6.95, 2.1, 5.6, 4.5, "Live Demo Sequence", [
    "Create user and auto-generated account",
    "Open second account (savings/checking)",
    "Deposit, withdraw, and transfer funds",
    "Show dashboard refresh and transaction logs",
    "Filter transaction history by date and type",
])
add_footer(slide)

# Slide 11: Risks and Limitations
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Current Limitations and Risks")
add_bullets(slide, 0.9, 2.0, 12.0, 4.9, [
    "Single-app deployment model: no horizontal scaling strategy yet.",
    "No automated test suite currently included in repository.",
    "No role-based permissions beyond per-user account scoping.",
    "Limited observability (metrics, tracing, and structured logs).",
    "No external integrations yet (payments, notifications, statement export).",
], font_size=21)
add_footer(slide)

# Slide 12: Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide)
add_title(slide, "Roadmap and Next Milestones")
add_card(slide, 0.8, 2.0, 3.8, 4.8, "Phase 1", [
    "Automated tests",
    "CSRF + security hardening",
    "Improved validation",
])
add_card(slide, 4.85, 2.0, 3.8, 4.8, "Phase 2", [
    "Admin and reporting panel",
    "Monthly statements (PDF/CSV)",
    "Notification engine",
])
add_card(slide, 8.9, 2.0, 3.8, 4.8, "Phase 3", [
    "Containerized deployment",
    "CI/CD and monitoring",
    "API layer for mobile clients",
])
add_footer(slide, "NexaBank Project Draft | Thank You")

output_file = "NexaBank_Project_Presentation.pptx"
prs.save(output_file)
print(f"Created: {output_file}")
